import streamlit as st
from modules.summarization import get_gemini_response
from modules.pdf_generator import generate_pdf_of_youtube_summaries
from modules.db_utils import add_to_db
from modules.youtube_utils import fetch_transcript
from modules.data_extraction import extract_numerical_data
from modules.ask_questions import ask_question, write_conversation_history
from modules.timeline_generator import extract_timeline
import requests
from bs4 import BeautifulSoup
import pyperclip
import io
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile
import os

def is_youtube_url(url):
    return url and ("youtube.com/watch" in url or "youtu.be" in url)


def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text


def extract_text_from_docx(file):
    doc = Document(io.BytesIO(file.read()))
    return "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])


def extract_text_from_pptx(file):
    prs = Presentation(io.BytesIO(file.read()))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def process_uploaded_files(uploaded_files):
    all_text = []
    for uploaded_file in uploaded_files:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        try:
            if file_extension == 'pdf':
                text = extract_text_from_pdf(io.BytesIO(uploaded_file.getvalue()))
            elif file_extension in ['doc', 'docx']:
                text = extract_text_from_docx(uploaded_file)
            elif file_extension in ['ppt', 'pptx']:
                text = extract_text_from_pptx(uploaded_file)
            else:
                continue
            all_text.append(f"=== Content from {uploaded_file.name} ===\n{text}")
        except Exception as e:
            st.error(f"Error processing {uploaded_file.name}: {str(e)}")
    return "\n\n".join(all_text)


def MainPage():
    st.title("YouTube Video Analysis & Quiz Generator")
    
    # File Upload Section at the top
    st.subheader("Upload Documents (Optional)")
    uploaded_files = st.file_uploader("Upload PDF, DOC, DOCX, PPT, or PPTX files", 
                                   type=['pdf', 'doc', 'docx', 'ppt', 'pptx'],
                                   accept_multiple_files=True,
                                   help="Upload documents to include in the analysis")
    if uploaded_files:
        st.session_state['uploaded_files'] = uploaded_files
    else:
        st.session_state['uploaded_files'] = []
    
    # Initialize session state variables
    if 'uploaded_files' not in st.session_state:
        st.session_state['uploaded_files'] = []
    if 'youtube_urls' not in st.session_state:
        st.session_state['youtube_urls'] = []
    if 'transcripts' not in st.session_state:
        st.session_state['transcripts'] = []
    if 'summary' not in st.session_state:
        st.session_state['summary'] = ""
    if 'combined_transcripts' not in st.session_state:
        st.session_state['combined_transcripts'] = ""
    if 'conversation_history' not in st.session_state:
        st.session_state['conversation_history'] = []
    if 'insights' not in st.session_state:
        st.session_state['insights'] = ""
    if 'numerical_data' not in st.session_state:
        st.session_state['numerical_data'] = ""
    if 'extracted_timeline' not in st.session_state:
        st.session_state['extracted_timeline'] = ""
    if 'fetch_summary_clicked' not in st.session_state:
        st.session_state['fetch_summary_clicked'] = False
    if 'collect_insights_clicked' not in st.session_state:
        st.session_state['collect_insights_clicked'] = False

    col1, col2 = st.columns([4, 1])
    with col1:
        input_url = st.text_input("Enter a URL (YouTube or Website)", placeholder="Enter the URL here...")
    with col2:
        st.markdown("<div style='height: 1.7em;'></div>", unsafe_allow_html=True)
        if st.button("Add", key="add_url_button"):
            if input_url.strip() == "":
                st.error("Please enter a valid URL.")
            elif input_url in st.session_state['youtube_urls']:
                st.warning("This URL is already added.")
            else:
                st.session_state['youtube_urls'].append(input_url)

    # Display list of added URLs with remove buttons
    if st.session_state['youtube_urls'] or (st.session_state.get('uploaded_files') and not st.session_state.get('fetch_summary_clicked', False)):
        st.subheader("📚 Added URLs")
        for idx, url in enumerate(st.session_state['youtube_urls']):
            col1, col2 = st.columns([8, 1])
            with col1:
                st.markdown(f"- {url}")
            with col2:
                if st.button("❌", key=f"remove_{idx}"):
                    st.session_state['youtube_urls'].pop(idx)
                    st.rerun()

    st.subheader("Fetch Transcripts and Generate Summary")
    
    # Add secondary prompt text box
    secondary_prompt = st.text_area("Additional Instructions (Optional)", 
                                  placeholder="Add any specific instructions for the summary...",
                                  height=70)
    
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("Fetch Summary", key="fetch_summary_button"):
            st.session_state['fetch_summary_clicked'] = True
            st.session_state['collect_insights_clicked'] = False
    with col2:
        if st.button("Collect Insights", key="collect_insights_button"):
            st.session_state['collect_insights_clicked'] = True
            st.session_state['fetch_summary_clicked'] = False
    
    fetch_summary_clicked = st.session_state['fetch_summary_clicked']
    collect_insights_clicked = st.session_state['collect_insights_clicked']

    if st.session_state.get('fetch_summary_clicked'):
        if not st.session_state['youtube_urls'] and not st.session_state.get('uploaded_files'):
            st.error("Please add at least one YouTube URL or upload a document.")
            st.session_state['fetch_summary_clicked'] = False
        else:
            # Process YouTube URLs first
            st.session_state['transcripts'] = []
            if st.session_state['youtube_urls']:
                for url in st.session_state['youtube_urls']:
                    if is_youtube_url(url):
                        transcript_content, error = fetch_transcript(url)
                        if error:
                            st.error(f"Error fetching transcript for {url}: {error}")
                        else:
                            st.session_state['transcripts'].append(transcript_content)
                    else:
                        try:
                            response = requests.get(url, timeout=10)
                            response.raise_for_status()
                            soup = BeautifulSoup(response.text, "html.parser")
                            for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                                tag.decompose()
                            text = soup.get_text(separator="\n")
                            clean_lines = [line.strip() for line in text.splitlines() if line.strip()]
                            page_content = "\n".join(clean_lines)
                            st.session_state['transcripts'].append(page_content[:10000])
                        except Exception as e:
                            st.error(f"Failed to fetch website content from {url}: {e}")
            
            # Process uploaded files
            documents_content = ""
            if st.session_state.get('uploaded_files'):
                with st.spinner("Processing uploaded files..."):
                    documents_content = process_uploaded_files(st.session_state['uploaded_files'])
            
            # Combine all content
            youtube_content = "\n\n".join(st.session_state['transcripts']) if st.session_state['transcripts'] else ""
            combined_content = "\n\n".join(filter(None, [youtube_content, documents_content]))
            
            if not combined_content.strip():
                st.error("No content available to summarize. Please check your inputs.")
                st.session_state['fetch_summary_clicked'] = False
                return
                
            with st.spinner('Generating summary from Gemini...'):
                base_prompt = '''Summarize the content as headers and paragraphs. Cover all the topics in 5 lines each. Do not miss even a single topic. Don't overuse bullet points. Use them only for important facts and numbers. '''
                # Append secondary prompt if provided
                if secondary_prompt and secondary_prompt.strip():
                    base_prompt = f"{base_prompt} \n\nAdditional instructions: {secondary_prompt}"
                
                gemini_response = get_gemini_response(combined_content + "\n\n" + base_prompt, model_name="gemini-2.0-flash")
                st.session_state['summary'] = gemini_response
                st.session_state['combined_transcripts'] = combined_content
                st.session_state['fetch_summary_clicked'] = False  # Reset the button state

    if st.session_state.get('collect_insights_clicked'):
        if not st.session_state['youtube_urls'] and not st.session_state.get('uploaded_files'):
            st.error("Please add at least one YouTube URL or upload a document.")
            st.session_state['collect_insights_clicked'] = False
        else:
            st.session_state['transcripts'] = []
            for url in st.session_state['youtube_urls']:
                if is_youtube_url(url):
                    transcript_content, error = fetch_transcript(url)
                    if error:
                        st.error(f"Error fetching transcript for {url}: {error}")
                    else:
                        st.session_state['transcripts'].append(transcript_content)
                else:
                    try:
                        response = requests.get(url, timeout=10)
                        response.raise_for_status()
                        soup = BeautifulSoup(response.text, "html.parser")
                        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
                            tag.decompose()
                        text = soup.get_text(separator="\n")
                        clean_lines = [line.strip() for line in text.splitlines() if line.strip()]
                        page_content = "\n".join(clean_lines)
                        st.session_state['transcripts'].append(page_content[:10000])
                    except Exception as e:
                        st.error(f"Failed to fetch website content from {url}: {e}")
            combined_transcripts = "\n\n".join(st.session_state['transcripts'])
            insights_prompt = '''I have some transcripts of news with me. Collect me some insights on them.
I don't want plain summaries. I want you to go deep, find patterns across multiple newsletters, and give me key emerging trends you notice.

✅ Avoid generic, overused statements like "AI is changing the world." Instead, explain specifically how such trends are unfolding—whether through new business models, shifts in user behavior, regulatory changes, or technological advancements.

✅ Structure your output clearly:

Start with a TL;DR section summarizing the key insights in 4–6 sentences.

Then go into detailed analysis using headers and paragraphs.

Use bullet points only when citing specific facts, numbers, or data points, not for general narration.

End with a summary of the key trends for quick reference.

Be analytical, connect the dots across sectors, and highlight what's genuinely new or noteworthy—not what's obvious or widely known.

Transcript : ''' + combined_transcripts
            with st.spinner('Collecting insights from Gemini...'):
                insights_response = get_gemini_response(insights_prompt, model_name="gemini-2.0-flash")
            st.session_state['insights'] = insights_response
            st.session_state['combined_transcripts'] = combined_transcripts

    if st.session_state.get('summary'):
        # Add Copy Transcripts button
        if st.button("📋 Copy Transcripts", key="copy_transcripts_button", 
                    help="Copy all transcripts to clipboard"):
            if 'combined_transcripts' in st.session_state and st.session_state['combined_transcripts']:
                pyperclip.copy(st.session_state['combined_transcripts'])
                st.success("Transcripts copied to clipboard!")
            else:
                st.warning("No transcripts available to copy.")
        
        st.write(st.session_state['summary'])

    if st.session_state.get('insights'):
        st.markdown("---")
        st.subheader(":bulb: Insights")
        st.write(st.session_state['insights'])

    # Extract Numerical Data    

    if st.button("Extract Numbers", key="extract_numericals"):
        with st.spinner("Extracting Numerical Data"):
            st.session_state['numerical_data'] = extract_numerical_data(st.session_state["summary"])
    if st.session_state['numerical_data']: st.markdown(st.session_state['numerical_data'])
    
    # Generate Timeline

    if st.button("Generate Timeline", key="generate_timeline_button"):
        with st.spinner("Generating timeline..."):
            st.session_state['extracted_timeline'] = extract_timeline(st.session_state["summary"])
    if st.session_state['extracted_timeline']: st.markdown(st.session_state['extracted_timeline'])
    
    # Ask Questions about Summary
    
    st.subheader("Ask Questions")
    if st.session_state['conversation_history']:
            write_conversation_history()
    question = st.text_input("Your Question", placeholder="Ask something about the transcript...")
    if st.button("Get Answer"):
        answer = ask_question(question, st.session_state['combined_transcripts']) 
        
        

    # This section has been moved and integrated into the main fetch_summary_clicked block above

    # Add Summary to DB
    if st.button("Save Summary to DB"):
        
        pdf_file = generate_pdf_of_youtube_summaries()

        # ✅ Run filldb.py to process the PDF into ChromaDB
        add_to_db(pdf_file)